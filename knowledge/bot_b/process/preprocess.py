#!/usr/bin/env python3
"""
Bot B 知识库预处理脚本
将蓝图(docx)和操作手册(pptx)转换为结构化知识单元

处理策略:
  蓝图(docx): 按 Heading 层级切分，段落+表格合并，保留标题路径作上下文
              跳过需求跟踪表(Business Requirement List)，拆分超大表格
  操作手册(pptx): 按 Agenda 分节，合并连续同功能步骤，最少 250 字

Usage:
    python3 preprocess.py [--all] [--blueprints] [--manuals]
"""

import os
import re
import json
import argparse
from pathlib import Path
from collections import defaultdict

from docx import Document
from pptx import Presentation

# ============================================================================
# 配置
# ============================================================================

BOT_B_DIR = Path(__file__).parent.parent
RAW_DIR = BOT_B_DIR / "raw"
PROCESSED_DIR = BOT_B_DIR / "processed"

CHUNK_MIN_CHARS = 150       # 低于此阈值直接丢弃
MANUAL_MERGE_MIN = 250      # 操作手册合并后最低字数
CHUNK_TARGET_MAX = 1500     # 超过此值尝试拆分
CHUNK_HARD_MAX = 2000       # 绝对上限

MODULE_MAP = {
    "POS": "POS",
    "Inventory": "Inventory",
    "Master data": "Master Data",
    "Promotion": "Promotion",
    "Data Dashboard": "Data Dashboard",
    "System Integration": "System Integration",
}


def detect_module(filename: str) -> str:
    for keyword, module in MODULE_MAP.items():
        if keyword.lower() in filename.lower():
            return module
    return "Unknown"


# ============================================================================
# 通用工具
# ============================================================================

def split_text(text: str, max_len: int, min_first: int = 200) -> list[str]:
    """Split text at paragraph boundaries, respecting max_len"""
    if len(text) <= max_len:
        return [text]

    parts = []
    remaining = text
    while remaining:
        if len(remaining) <= max_len:
            parts.append(remaining)
            break
        # Try double newline first, then single newline
        cut = remaining.rfind("\n\n", min_first, max_len)
        if cut == -1:
            cut = remaining.rfind("\n", min_first, max_len)
        if cut == -1:
            cut = max_len
        parts.append(remaining[:cut].strip())
        remaining = remaining[cut:].strip()
    return parts


# ============================================================================
# 蓝图(docx)处理
# ============================================================================

def is_requirement_tracking_table(table) -> bool:
    """Detect project tracking tables (low knowledge value, huge size)"""
    if len(table.rows) < 3:
        return False
    header = " ".join(cell.text.strip().lower() for cell in table.rows[0].cells)
    tracking_keywords = ["as-is", "to-be", "not in rfp", "chapters involved", "no change"]
    matches = sum(1 for kw in tracking_keywords if kw in header)
    return matches >= 3


def table_to_text(table) -> str:
    """Convert table to readable text, handling 2-column vs multi-column layouts"""
    if len(table.rows) < 2:
        return ""

    num_cols = max(len(row.cells) for row in table.rows)

    # For wide tracking tables, skip them
    if is_requirement_tracking_table(table):
        return ""

    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        cells = [re.sub(r"\s+", " ", c) for c in cells]
        # Deduplicate merged cells (python-docx repeats content)
        deduped = []
        for c in cells:
            if c and (not deduped or c != deduped[-1]):
                deduped.append(c)
        if deduped:
            rows.append(" | ".join(deduped))

    return "\n".join(rows)


def extract_blueprint_chunks(filepath: Path) -> list[dict]:
    doc = Document(filepath)
    module = detect_module(filepath.name)
    source = filepath.name

    heading_stack = []
    chunks = []
    current_parts = []   # accumulated text parts (paragraphs + tables)
    current_size = 0

    def flush():
        nonlocal current_parts, current_size
        if not current_parts:
            return
        content = "\n\n".join(current_parts).strip()
        current_parts, current_size = [], 0
        if len(content) < CHUNK_MIN_CHARS:
            return

        context = " > ".join(heading_stack) if heading_stack else module

        # Split if too large
        for part in split_text(content, CHUNK_HARD_MAX):
            if len(part) >= CHUNK_MIN_CHARS:
                chunks.append({
                    "content": part,
                    "context": context,
                    "module": module,
                    "source_file": source,
                    "doc_type": "blueprint",
                    "heading_path": list(heading_stack),
                })

    from docx.oxml.ns import qn

    body = doc.element.body
    table_idx = 0

    for child in body:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            para = None
            for p in doc.paragraphs:
                if p._element is child:
                    para = p
                    break
            if para is None:
                continue

            text = para.text.strip()
            if not text:
                continue

            style = para.style.name if para.style else ""

            if style.startswith("Heading"):
                try:
                    level = int(style.replace("Heading ", "").strip())
                except ValueError:
                    level = 1

                # Flush accumulated content before heading change
                flush()

                heading_stack = heading_stack[: level - 1]
                heading_stack.append(text)
            else:
                current_parts.append(text)
                current_size += len(text)
                if current_size > CHUNK_HARD_MAX * 1.5:
                    flush()

        elif tag == "tbl":
            if table_idx < len(doc.tables):
                table = doc.tables[table_idx]
                table_idx += 1
                ttext = table_to_text(table)
                if not ttext or len(ttext) < 20:
                    continue
                current_parts.append(ttext)
                current_size += len(ttext)
                # Tables can be very large, flush immediately
                if current_size > CHUNK_HARD_MAX:
                    flush()

    flush()
    return chunks


def process_blueprints() -> list[dict]:
    print("\n[1/2] 处理蓝图文档 (docx)...")
    all_chunks = []

    for fname in sorted(RAW_DIR.glob("*.docx")):
        if fname.name.startswith("~") or fname.name.startswith("."):
            continue
        chunks = extract_blueprint_chunks(fname)
        print(f"  {fname.name}: {len(chunks)} 个知识单元")
        all_chunks.extend(chunks)

    print(f"  蓝图合计: {len(all_chunks)} 个知识单元")
    return all_chunks


# ============================================================================
# 操作手册(pptx)处理
# ============================================================================

def is_agenda_slide(slide) -> bool:
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text.strip().lower())
    combined = " ".join(texts).strip()
    return combined == "agenda" or (combined.startswith("agenda") and len(combined) < 20)


def extract_slide_text(slide) -> tuple[str, str]:
    title = ""
    body_parts = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text.strip()
        if not full_text or len(full_text) <= 2:
            continue
        if not title and len(full_text) > 3:
            title = full_text.split("\n")[0][:120]
        body_parts.append(full_text)

    return title, "\n".join(body_parts).strip()


def extract_operation_name(title: str, body: str) -> str:
    """Extract operation name from slide title or body"""
    if title:
        # Remove leading page numbers like "42\n" or just "42"
        cleaned = re.sub(r"^\d+[\s\n]*", "", title).strip()
        # Skip generic labels
        if cleaned and cleaned.lower() not in ("operation steps:", "steps/explanation", "agenda"):
            return cleaned[:80]
    return ""


def extract_manual_chunks(filepath: Path) -> list[dict]:
    prs = Presentation(filepath)
    module = detect_module(filepath.name)
    source = filepath.name

    # Collect all slides with metadata
    slides_data = []
    for i, slide in enumerate(prs.slides):
        agenda = is_agenda_slide(slide)
        title, body = extract_slide_text(slide)
        op_name = extract_operation_name(title, body) if not agenda else ""
        slides_data.append({
            "index": i,
            "is_agenda": agenda,
            "title": title,
            "body": body,
            "body_len": len(body),
            "op_name": op_name,
        })

    # Group slides into sections by Agenda dividers
    sections = []
    current_name = f"{module}"
    current_slides = []

    for sd in slides_data:
        if sd["is_agenda"]:
            if current_slides:
                sections.append((current_name, current_slides))
            current_name = module
            current_slides = []
            continue
        if sd["body_len"] > 15:
            if not current_slides and sd["op_name"]:
                current_name = f"{module} - {sd['op_name']}"
            current_slides.append(sd)

    if current_slides:
        sections.append((current_name, current_slides))

    # Convert sections to chunks with merging
    chunks = []
    for section_name, section_slides in sections:
        # Merge consecutive slides into operation groups
        # Use operation name changes as boundaries
        raw_groups = []
        current_op = ""
        current_bodies = []

        for sd in section_slides:
            body = sd["body"]
            op = sd["op_name"]

            # New operation if name changes and we have content
            if op and op != current_op and current_bodies:
                raw_groups.append((current_op, current_bodies))
                current_bodies = []

            if op:
                current_op = op
            current_bodies.append(body)

        if current_bodies:
            raw_groups.append((current_op, current_bodies))

        # Now merge small groups into chunks of target size
        pending_text = ""
        pending_op = ""

        for op_name, bodies in raw_groups:
            combined = "\n\n".join(b for b in bodies if len(b) > 15)
            if not combined.strip():
                continue

            if pending_text:
                candidate = pending_text + "\n\n" + combined
                if len(candidate) <= CHUNK_TARGET_MAX:
                    pending_text = candidate
                    continue
                else:
                    # Flush pending
                    if len(pending_text) >= MANUAL_MERGE_MIN:
                        for part in split_text(pending_text, CHUNK_HARD_MAX):
                            chunks.append({
                                "content": part,
                                "context": pending_op or section_name,
                                "module": module,
                                "source_file": source,
                                "doc_type": "user_manual",
                                "operation_name": pending_op or section_name,
                            })
                    pending_text = combined
                    pending_op = op_name or pending_op
            else:
                pending_text = combined
                pending_op = op_name or section_name

        # Flush remaining
        if pending_text and len(pending_text) >= MANUAL_MERGE_MIN:
            for part in split_text(pending_text, CHUNK_HARD_MAX):
                chunks.append({
                    "content": part,
                    "context": pending_op or section_name,
                    "module": module,
                    "source_file": source,
                    "doc_type": "user_manual",
                    "operation_name": pending_op or section_name,
                })

    return chunks


def process_manuals() -> list[dict]:
    print("\n[2/2] 处理操作手册 (pptx)...")
    all_chunks = []

    for fname in sorted(RAW_DIR.glob("*.pptx")):
        if fname.name.startswith("~") or fname.name.startswith("."):
            continue
        chunks = extract_manual_chunks(fname)
        print(f"  {fname.name}: {len(chunks)} 个知识单元")
        all_chunks.extend(chunks)

    print(f"  操作手册合计: {len(all_chunks)} 个知识单元")
    return all_chunks


# ============================================================================
# 输出生成
# ============================================================================

def assign_doc_ids(chunks: list[dict]) -> list[dict]:
    for i, chunk in enumerate(chunks):
        module = chunk.get("module", "unknown").lower().replace(" ", "_")
        doc_type = chunk.get("doc_type", "unknown")
        chunk["doc_id"] = f"botb_{doc_type}_{module}_{str(i).zfill(4)}"
    return chunks


def generate_json(chunks: list[dict], output_path: Path):
    output = []
    for c in chunks:
        output.append({
            "doc_id": c["doc_id"],
            "title": c.get("context", ""),
            "content": c["content"],
            "metadata": {
                "module": c.get("module", ""),
                "doc_type": c.get("doc_type", ""),
                "source_file": c.get("source_file", ""),
                "heading_path": c.get("heading_path", []),
                "operation_name": c.get("operation_name", ""),
            },
        })
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output, f, ensure_ascii=False, indent=2)


def generate_txt_for_dify(chunks: list[dict], output_dir: Path):
    groups = defaultdict(list)
    for c in chunks:
        module = c.get("module", "unknown").lower().replace(" ", "_")
        doc_type = c.get("doc_type", "unknown")
        key = f"{doc_type}_{module}"
        groups[key].append(c)

    output_dir.mkdir(parents=True, exist_ok=True)
    file_list = []

    for key, group_chunks in sorted(groups.items()):
        batch_size = 300
        for batch_idx in range(0, len(group_chunks), batch_size):
            batch = group_chunks[batch_idx: batch_idx + batch_size]
            fname = f"{key}"
            if batch_idx > 0:
                fname += f"_part{batch_idx // batch_size + 1}"
            fname += ".txt"

            lines = []
            for c in batch:
                lines.append("=" * 60)
                context = c.get("context", "")
                if context:
                    lines.append(f"[{context}]")
                lines.append("")
                lines.append(c["content"])
                lines.append("")

            filepath = output_dir / fname
            with open(filepath, "w", encoding="utf-8") as f:
                f.write("\n".join(lines))

            size_kb = filepath.stat().st_size / 1024
            file_list.append((fname, len(batch), size_kb))

    return file_list


def print_stats(chunks: list[dict]):
    print("\n" + "=" * 70)
    print("知识单元统计")
    print("=" * 70)

    by_type = defaultdict(list)
    by_module = defaultdict(list)
    lengths = []

    for c in chunks:
        by_type[c.get("doc_type", "?")].append(c)
        by_module[c.get("module", "?")].append(c)
        lengths.append(len(c["content"]))

    print("\n按文档类型:")
    for dt, items in sorted(by_type.items()):
        lens = [len(c["content"]) for c in items]
        print(f"  {dt:15s}: {len(items):>4} 个  "
              f"平均 {sum(lens)//len(lens):>5} 字  "
              f"最小 {min(lens):>5}  最大 {max(lens):>5}")

    print("\n按业务模块:")
    for mod, items in sorted(by_module.items()):
        lens = [len(c["content"]) for c in items]
        print(f"  {mod:20s}: {len(items):>4} 个  "
              f"平均 {sum(lens)//len(lens):>5} 字  "
              f"总计 {sum(lens):>8,} 字")

    print(f"\n总体: {len(chunks)} 个知识单元, "
          f"总文本 {sum(lengths):,} 字, "
          f"平均 {sum(lengths)//len(lengths)} 字")

    ranges = [(0, 200), (200, 500), (500, 1000), (1000, 1500), (1500, 2000), (2000, float("inf"))]
    labels = ["<200", "200-500", "500-1000", "1000-1500", "1500-2000", ">2000"]
    print("\n长度分布:")
    for (lo, hi), label in zip(ranges, labels):
        count = sum(1 for l in lengths if lo <= l < hi)
        bar = "#" * (count // 5)
        print(f"  {label:>10s}: {count:>4} {bar}")


# ============================================================================
# 主函数
# ============================================================================

def main():
    parser = argparse.ArgumentParser(description="Bot B 知识库预处理")
    parser.add_argument("--all", action="store_true", help="处理所有文档")
    parser.add_argument("--blueprints", action="store_true", help="只处理蓝图(docx)")
    parser.add_argument("--manuals", action="store_true", help="只处理操作手册(pptx)")
    args = parser.parse_args()

    if not any([args.all, args.blueprints, args.manuals]):
        args.all = True

    PROCESSED_DIR.mkdir(parents=True, exist_ok=True)

    print("=" * 70)
    print("Bot B 知识库预处理")
    print("=" * 70)

    all_chunks = []

    if args.all or args.blueprints:
        bp_chunks = process_blueprints()
        all_chunks.extend(bp_chunks)

    if args.all or args.manuals:
        manual_chunks = process_manuals()
        all_chunks.extend(manual_chunks)

    if not all_chunks:
        print("\n未生成任何知识单元，请检查 raw/ 目录")
        return

    all_chunks = assign_doc_ids(all_chunks)

    json_path = PROCESSED_DIR / "knowledge_units.json"
    generate_json(all_chunks, json_path)
    print(f"\n已保存 JSON: {json_path} ({json_path.stat().st_size / 1024:.1f} KB)")

    txt_dir = PROCESSED_DIR / "dify_upload"
    txt_files = generate_txt_for_dify(all_chunks, txt_dir)
    print(f"\nDify 导入文件 ({len(txt_files)} 个):")
    for fname, count, size_kb in txt_files:
        print(f"  {fname:45s} {count:>4} 条  {size_kb:>8.1f} KB")

    print_stats(all_chunks)

    print("\n预处理完成!")


if __name__ == "__main__":
    main()
